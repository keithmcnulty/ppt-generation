from pptx import Presentation
from pptx.chart.data import CategoryChartData


def edit_pres(group, data_chart, data_table, input_file, output_file):
  
  # transform input data for insertion into ppt
  series = data_chart[data_chart["group"] == group].squeeze()
  sales_data = data_table
  
  # get presentation template
  pres = Presentation(input_file)

  # title slide
  slide0 = pres.slides[0]

  # change title
  title = [s for s in slide0.shapes if s.has_text_frame and s.text.find("Presentation title")!=-1]
  title[0].text = "Presentation for Group " + group

  # change subtitle
  subtitle = [s for s in slide0.shapes if s.has_text_frame and s.text.find("Subtitle")!=-1]
  subtitle[0].text = "Financial Information for Group " + group

  # first chart
  slide1 = pres.slides[1]

  # change title
  title = [s for s in slide1.shapes if s.has_text_frame and s.text.find("Chart")!=-1]
  title[0].text = "Financial Results Summary for Group " + group

  # get charts
  charts = [s for s in slide1.shapes if s.has_chart]

  # change bar chart
  chart0 = charts[0]
  chart_data = CategoryChartData()
  chart_data.categories = ['Category ' + str(i) for i in range(1,5)]
  chart_data.add_series('Series 1', series[['cat1_1', 'cat2_1', 'cat3_1', 'cat4_1']].values)
  chart_data.add_series('Series 2', series[['cat1_2', 'cat2_2', 'cat3_2', 'cat4_2']].values)
  chart_data.add_series('Series 3', series[['cat1_3', 'cat2_3', 'cat3_3', 'cat4_3']].values)
  chart0.chart.replace_data(chart_data)
  chart0.chart.chart_title.text_frame.text = "Sales by Category: Group " + group

  # change pie chart
  chart1 = charts[1]
  chart_data.categories = ["1st Qtr", "2nd Qtr", "3rd Qtr", "4th Qtr"]
  chart_data.add_series('Series 1', series[['pie1', 'pie2', 'pie3', 'pie4']].values)
  chart1.chart.replace_data(chart_data)
  chart1.chart.chart_title.text_frame.text = "Sales by Quarter: Group " + group

  # second chart
  slide2 = pres.slides[2]

  # change title
  title = [s for s in slide2.shapes if s.has_text_frame and s.text.find("Table")!=-1]
  title[0].text = "Results Table for Group " + group

  # replace table
  table = [s for s in slide2.shapes if s.has_table]

  # headings
  for j in range(0,5):
    table[0].table.cell(0, j).text = "Product " + sales_data.columns.values[j]

  # data
  for i in range(1,9):
    for j in range(0,5):
      table[0].table.cell(i,j).text = str(sales_data.iloc[i-1, j])

  # totals
  for j in range(0,5):
    table[0].table.cell(9,j).text = "{:.1f}".format(sum(sales_data.iloc[:,j]))
  
  # save new pres and generate success message
  pres.save(output_file)
  return "Successfully saved version " + group + "!"
